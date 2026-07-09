
import streamlit as st
import pandas as pd
from pathlib import Path
import re, io

st.set_page_config(page_title="LabSafeEdu参赛版", page_icon="🧪", layout="wide")
DATA_DIR = Path(__file__).parent / "data"

@st.cache_data
def load_data():
    return {
        "knowledge": pd.read_csv(DATA_DIR / "knowledge_base.csv"),
        "hazards": pd.read_csv(DATA_DIR / "hazard_rectification_library.csv"),
        "risk": pd.read_csv(DATA_DIR / "risk_rules.csv"),
        "cases": pd.read_csv(DATA_DIR / "scenario_cases.csv"),
    }

data = load_data()

def norm(x):
    return str(x).lower().replace("，"," ").replace("。"," ").replace("、"," ").replace("；"," ")

def has(q, words):
    q = norm(q)
    return any(w in q for w in words)

def risk_assess(text):
    q = norm(text)
    score, factors = 0, []
    for _, r in data["risk"].iterrows():
        kws = str(r.get("关键词", "")).split()
        if any(k.lower() in q for k in kws):
            points = int(r.get("分值", 1))
            score += points
            factors.append((r.get("风险因素", ""), points, r.get("控制措施", "")))
    level = "重大风险" if score >= 8 else ("高风险" if score >= 4 else ("中风险" if score >= 2 else "低风险"))
    return level, score, factors

def hazard_match(text):
    q = norm(text)
    best, score = None, 0
    for _, r in data["hazards"].iterrows():
        alltxt = " ".join(map(str, [r.get("隐患描述",""), r.get("隐患类别",""), r.get("问题分析",""), r.get("整改建议",""), r.get("关键词","")])).lower()
        tokens = re.findall(r"[\u4e00-\u9fff]{2,}|[a-zA-Z0-9]+", q)
        s = sum(1 for w in tokens if w in alltxt)
        if s > score:
            best, score = r, s
    if best is not None and score > 0:
        return best, "隐患库匹配"
    if has(q, ["废液", "矿泉水", "饮料瓶", "无标签"]):
        return pd.Series({"隐患类别":"实验废弃物管理","风险等级":"高风险","问题分析":"废液容器或标签不规范，可能造成成分不明、误用、泄漏或混装反应。","整改建议":"立即使用专用废液桶分类收集，补充标签并完善台账。","复查要点":"复查容器、标签、成分、责任人和暂存区域。"}), "智能兜底"
    if has(q, ["气瓶", "钢瓶"]):
        return pd.Series({"隐患类别":"气瓶安全","风险等级":"高风险","问题分析":"气瓶管理不规范可能造成倾倒、泄漏、窒息或燃爆风险。","整改建议":"固定气瓶，检查阀门、减压阀、压力表和管路连接。","复查要点":"复查固定状态、阀门、管路、检漏记录和台账。"}), "智能兜底"
    if has(q, ["护目镜", "手套", "实验服", "防护"]):
        return pd.Series({"隐患类别":"个人防护","风险等级":"中风险","问题分析":"个人防护不足可能造成飞溅、腐蚀、割伤或吸入暴露。","整改建议":"停止不规范操作，按风险佩戴合适PPE并开展再教育。","复查要点":"复查PPE佩戴、现场告知和应急设施可达性。"}), "智能兜底"
    return pd.Series({"隐患类别":"一般安全隐患","风险等级":"中风险","问题分析":"现场行为可能影响安全秩序、风险控制或责任追溯。","整改建议":"记录问题，明确责任人、整改时限并复查销号。","复查要点":"复查整改记录、责任人、完成情况和教育反馈。"}), "智能兜底"

def build_access_plan(identity, lab_type, task):
    q = norm(lab_type + task)
    focus = ["准入", "防护", "应急"]
    if has(q, ["化学", "材料", "试剂", "溶剂", "废液"]):
        focus += ["危化品", "废液", "气瓶", "用电"]
    df = data["knowledge"]
    selected = df[df.apply(lambda r: any(k in str(r.to_dict()) for k in focus), axis=1)]
    if selected.empty:
        selected = df
    if identity in ["研究生", "本科生"]:
        note = "重点强化实验前风险辨识、PPE佩戴、废液分类、气瓶/用电安全和离室检查。"
    elif identity in ["新进教师", "实验技术人员"]:
        note = "重点强化责任落实、风险分级、台账管理和隐患闭环。"
    else:
        note = "重点完成安全告知、陪同进入、禁止独立操作和现场设施确认。"
    return selected.head(12), note

def parse_upload(uploaded_file, extra_text=""):
    text = extra_text.strip()
    notes = []
    if uploaded_file is None:
        return text, notes
    name = uploaded_file.name.lower()
    b = uploaded_file.getvalue()
    try:
        if name.endswith(".csv"):
            df = pd.read_csv(io.BytesIO(b))
            text += "\n" + df.to_csv(index=False)
        elif name.endswith((".xlsx", ".xls")):
            sheets = pd.read_excel(io.BytesIO(b), sheet_name=None)
            for s, df in sheets.items():
                text += f"\n【{s}】\n" + df.to_csv(index=False)
        elif name.endswith(".docx"):
            from docx import Document
            doc = Document(io.BytesIO(b))
            text += "\n" + "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        elif name.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(b))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                t = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        t.append(shape.text.strip())
                if t:
                    parts.append(f"【第{i}页】\n" + "\n".join(t))
            text += "\n" + "\n".join(parts)
        elif name.endswith(".pdf"):
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(b))
            text += "\n" + "\n".join([(p.extract_text() or "") for p in reader.pages])
        else:
            text += "\n" + b.decode("utf-8", errors="ignore")
    except Exception as e:
        notes.append(f"解析提示：{e}")
    return text.strip(), notes

st.sidebar.title("🧪 LabSafeEdu")
page = st.sidebar.radio("选择功能模块", ["首页","场景化准入学习","实验前风险研判","隐患整改闭环","案例微课堂","准入测评与错题反馈","管理数据看板","知识库维护"])
st.sidebar.caption("高校实验室安全准入与风险闭环教育智能体")

if page == "首页":
    st.title("LabSafeEdu实验室安全准入与风险闭环教育智能体")
    st.caption("围绕“场景化准入学习—实验前风险研判—隐患整改闭环—案例育人—测评反馈—管理优化”构建AI教育智能体")
    c1,c2,c3,c4 = st.columns(4)
    c1.metric("检查实验室", "133间")
    c2.metric("累计隐患", "832项")
    c3.metric("化学安全隐患", "384项")
    c4.metric("整改率", "89%→100%")
    st.subheader("报告—系统一致性")
    st.dataframe(pd.DataFrame([
        ["准入学习与真实实验场景脱节", "场景化准入学习", "按身份、实验室类型和实验任务生成个性化学习路径"],
        ["实验前风险研判不足", "实验前风险研判", "识别风险因素、生成风险等级、控制措施和检查清单"],
        ["隐患整改难以转化为教育资源", "隐患整改闭环 + 案例微课堂 + 测评反馈", "将巡检隐患转化为整改建议、复查要点、案例和测评题目"],
        ["制度与案例持续更新成本高", "知识库维护", "支持制度、培训材料和隐患台账导入，形成可迭代本地知识库"],
    ], columns=["拟解决的问题","对应功能","智能体作用"]), use_container_width=True)
    st.subheader("教育闭环")
    st.info("学前准入 → 实验前研判 → 隐患整改 → 案例育人 → 测评反馈 → 管理优化 → 知识库更新")

elif page == "场景化准入学习":
    st.header("场景化准入学习")
    identity = st.selectbox("学习对象", ["本科生","研究生","新进教师","实验技术人员","短期进入人员"])
    lab_type = st.selectbox("实验室类型", ["化学类实验室","材料类实验室","机电类实验室","生物类实验室","其他实验室"])
    task = st.text_area("拟进入/开展的实验场景", "研一新生第一次进入化学类实验室，后续可能开展有机溶剂实验。")
    if st.button("生成个性化准入学习路径"):
        selected, note = build_access_plan(identity, lab_type, task)
        st.success(note)
        st.dataframe(selected[[c for c in ["一级类别","二级类别","标题","内容"] if c in selected.columns]], use_container_width=True)
        st.info("推荐流程：制度学习 → 现场设施确认 → 准入考试 → 导师审核 → 首次实验陪同/复核。")

elif page == "实验前风险研判":
    st.header("实验前风险研判")
    text = st.text_area("输入实验内容", "今晚做有机溶剂加热回流实验，使用乙醇和丙酮，会产生有机废液。")
    col1, col2, col3 = st.columns(3)
    extra = []
    with col1:
        if st.checkbox("夜间实验"): extra.append("夜间")
        if st.checkbox("首次独立操作"): extra.append("首次 独立")
    with col2:
        if st.checkbox("使用气瓶"): extra.append("气瓶")
        if st.checkbox("高温/加热"): extra.append("加热 高温")
    with col3:
        if st.checkbox("大量有机溶剂"): extra.append("有机溶剂 大量")
        if st.checkbox("产生危险废液"): extra.append("废液 危废")
    if st.button("生成风险研判"):
        level, score, factors = risk_assess(text + " " + " ".join(extra))
        st.metric("风险等级", level)
        st.metric("风险分值", score)
        if factors:
            st.dataframe(pd.DataFrame(factors, columns=["风险因素","分值","控制措施"]), use_container_width=True)
        checklist = ["查阅SDS和实验方案","确认实验服、手套、护目镜等PPE","检查通风橱、冷凝水和加热设备","准备专用废液桶并张贴标签","明确导师/负责人和应急联系人","实验结束完成水、电、气、门窗离室检查"]
        if level in ["高风险", "重大风险"]:
            checklist += ["完成导师审核或实验室负责人审批","涉及夜间实验时完成夜间备案","不建议单人开展高风险实验"]
        st.dataframe(pd.DataFrame({"实验前检查清单": checklist, "完成情况": ["□"]*len(checklist)}), use_container_width=True)
        st.info("系统同步推荐：进入案例微课堂学习相关隐患案例，并完成废液、PPE、用电和应急类测评题目。")

elif page == "隐患整改闭环":
    st.header("隐患整改闭环")
    text = st.text_area("输入巡检隐患", "巡检发现学生用矿泉水桶装实验废液，桶身没有规范标签。")
    if st.button("生成整改闭环建议"):
        h, src = hazard_match(text)
        st.metric("风险等级", h["风险等级"])
        rows = [
            ("识别来源", src),
            ("隐患类别", h["隐患类别"]),
            ("问题分析", h["问题分析"]),
            ("整改建议", h["整改建议"]),
            ("复查要点", h["复查要点"]),
            ("教育转化", "可转化为案例微课堂、测评题目和后续培训重点。")
        ]
        st.dataframe(pd.DataFrame(rows, columns=["项目","内容"]), use_container_width=True)

elif page == "案例微课堂":
    st.header("案例微课堂")
    name = st.selectbox("选择案例", data["cases"]["案例名称"].tolist())
    row = data["cases"][data["cases"]["案例名称"] == name].iloc[0]
    st.metric("风险等级", row["风险等级"])
    for lab in ["案例描述","风险分析","正确做法","追问问题","教学目标"]:
        if lab in row.index:
            st.markdown(f"**{lab}：** {row[lab]}")

elif page == "准入测评与错题反馈":
    st.header("准入测评与错题反馈")
    st.write("围绕准入学习和典型隐患生成测评题，错题可反向推荐制度知识和案例微课堂内容。")
    questions = pd.DataFrame([
        ["单选", "实验废液应如何处理？", "分类收集于专用废液桶并贴标签", "废液管理"],
        ["判断", "夜间开展高风险实验前应进行备案。", "正确", "夜间实验"],
        ["情景", "巡检发现废液桶无标签，应如何整改？", "补贴标签、完善台账、复查销号", "隐患整改"],
        ["情景", "准备做有机溶剂加热回流实验，实验前应检查什么？", "SDS、PPE、通风橱、冷凝水、废液桶、备案和值守", "实验前风险研判"],
    ], columns=["题型","题目","答案/解析","知识点"])
    st.dataframe(questions, use_container_width=True)
    weak = st.multiselect("选择薄弱知识点生成复习建议", ["废液管理","夜间实验","隐患整改","个人防护","气瓶安全","实验前风险研判"])
    if weak:
        st.success("复习建议：进入制度知识库和案例微课堂学习：" + "、".join(weak))

elif page == "管理数据看板":
    st.header("管理数据看板")
    c1,c2,c3,c4 = st.columns(4)
    c1.metric("检查实验室", "133间")
    c2.metric("累计隐患", "832项")
    c3.metric("化学安全隐患", "384项")
    c4.metric("整改率", "89%→100%")
    df = pd.DataFrame([["实验场所方面",64],["安全设施方面",78],["基础安全方面",206],["化学安全方面",384],["生物安全方面",4],["其他隐患",96]], columns=["隐患类别","数量"])
    st.dataframe(df, use_container_width=True)
    st.bar_chart(df.set_index("隐患类别"), height=360)
    st.info("管理建议：以化学安全和基础安全为重点，强化危化品、废液、气瓶、用电、夜间实验和大型仪器连续运行管理。")

elif page == "知识库维护":
    st.header("知识库维护")
    st.write("本模块用于体现制度和案例的持续更新能力。管理者可上传制度、PPT培训材料、Excel隐患台账、PDF检查标准或文本资料，由系统生成入库建议。")
    uploaded = st.file_uploader("上传制度/台账/培训材料", type=["txt","md","csv","xlsx","xls","docx","pptx","pdf"])
    extra = st.text_area("或粘贴新增制度条款/隐患案例/检查标准")
    if st.button("解析并生成入库建议"):
        text, notes = parse_upload(uploaded, extra)
        for n in notes:
            st.warning(n)
        if not text:
            st.warning("请先上传文件或粘贴内容。")
        else:
            preview = text[:800]
            st.success("已解析资料，并生成入库建议。")
            st.text_area("资料内容预览", preview, height=180)
            q = norm(text)
            target = []
            if has(q, ["隐患", "整改", "巡检", "复查"]): target.append("隐患整改库/案例微课堂")
            if has(q, ["培训", "制度", "要求", "规范", "检查标准"]): target.append("制度知识库")
            if has(q, ["题", "考试", "测评", "选择", "判断"]): target.append("准入测评题库")
            if not target: target.append("制度知识库")
            st.dataframe(pd.DataFrame([
                ["建议入库模块", "、".join(target)],
                ["建议状态", "待管理者审核"],
                ["后续作用", "更新制度知识、隐患规则、案例微课堂和测评反馈内容"]
            ], columns=["项目","内容"]), use_container_width=True)
